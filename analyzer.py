"""
analyzer.py — Полный анализ .pptx презентаций
Типографика · Вёрстка · Шрифты · Интервалы · Буллиты · Изображения · Цвета · Таблицы
"""

import re
import zipfile
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from collections import defaultdict, Counter

# ─────────────────────────────────────────────
# КОНСТАНТЫ
# ─────────────────────────────────────────────

POSITION_TOLERANCE   = 182880    # ~0.2 см
SIZE_TOLERANCE       = 457200    # ~0.5 см
FONT_SIZE_TOLERANCE  = 2         # pt
GRID_BUCKET_CM       = 5.0       # бакет для дедупликации фигур

MIN_BODY_FONT_PT     = 14        # меньше — слишком мелко
MIN_TITLE_FONT_PT    = 20        # меньше — заголовок не заметен
MAX_FONT_SIZES_SLIDE = 4         # больше → пёстро
MAX_TITLE_LEN        = 80        # символов в заголовке
LINE_SPACING_MIN     = 0.8
LINE_SPACING_MAX     = 2.5
MAX_COLORS_SLIDE     = 6         # уникальных цветов на слайде
MAX_TABLE_ROWS       = 10
MAX_TABLE_COLS       = 7
OVERLAP_MIN_CM       = 0.5       # минимальное перекрытие, чтобы считать ошибкой
SNAP_THRESHOLD_CM    = 0.4       # «почти выровнено» — ближе чем это, но не совпадает

TITLE_KEYWORDS = ("title", "заголовок", "heading", "header")
NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


# ─────────────────────────────────────────────
# ВСПОМОГАТЕЛЬНЫЕ
# ─────────────────────────────────────────────

def emu_to_cm(emu):
    return round(emu / 914400 * 2.54, 2)

def get_shape_text(shape):
    if shape.has_text_frame:
        return "\n".join(p.text for p in shape.text_frame.paragraphs)
    return ""

def is_title_shape(shape):
    if any(kw in shape.name.lower() for kw in TITLE_KEYWORDS):
        return True
    try:
        ph = shape.placeholder_format
        if ph and ph.idx in (0, 1):
            return True
    except Exception:
        pass
    return False

def get_font_sizes(shape) -> set:
    sizes = set()
    if not shape.has_text_frame:
        return sizes
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                sizes.add(round(run.font.size / 12700, 1))
        try:
            if para.font.size:
                sizes.add(round(para.font.size / 12700, 1))
        except Exception:
            pass
    return sizes

def make_shape_key(name: str, left: int, top: int) -> str:
    l_b = round(emu_to_cm(left) / GRID_BUCKET_CM)
    t_b = round(emu_to_cm(top)  / GRID_BUCKET_CM)
    return f"{name}__L{l_b}__T{t_b}"

def get_bullet_type(para) -> str:
    pPr = para._p.pPr
    if pPr is None:
        return "inherit"
    if pPr.find(f"{{{NS}}}buNone") is not None:
        return "none"
    ch_el = pPr.find(f"{{{NS}}}buChar")
    if ch_el is not None:
        return f"char:{ch_el.get('char', '•')}"
    if pPr.find(f"{{{NS}}}buAutoNum") is not None:
        return "autonum"
    return "inherit"

def get_shape_colors(shape) -> set:
    colors = set()
    try:
        from pptx.enum.dml import MSO_THEME_COLOR
        fill = shape.fill
        if fill.type is not None:
            try:
                c = fill.fore_color.rgb
                colors.add(str(c))
            except Exception:
                pass
    except Exception:
        pass
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                try:
                    if run.font.color.type:
                        colors.add(str(run.font.color.rgb))
                except Exception:
                    pass
    return colors

def format_number(n_str: str) -> str:
    try:
        n = int(n_str)
        s = str(abs(n))
        parts = []
        for i, ch in enumerate(reversed(s)):
            if i > 0 and i % 3 == 0:
                parts.append(' ')
            parts.append(ch)
        return ('-' if n < 0 else '') + ''.join(reversed(parts))
    except Exception:
        return n_str


# ─────────────────────────────────────────────
# БЛОК 1: ТИПОГРАФИКА
# ─────────────────────────────────────────────

PREPOSITIONS = (
    "в","на","с","к","по","о","и","а","но","из","до","от","за","об",
    "под","над","при","про","для","без","или","не","то","бы","же","со","во","ко",
)
HANGING_RE      = re.compile(r'\b(' + '|'.join(PREPOSITIONS) + r')\s*$', re.IGNORECASE | re.UNICODE)
DOUBLE_SPACE    = re.compile(r'  +')
# Дефис вместо тире — ловит и двойные пробелы вокруг
HYPHEN_AS_DASH  = re.compile(
    r'(?<=[а-яёa-zА-ЯЁA-Z\d])\s+-\s+(?=[а-яёa-zА-ЯЁA-Z\d])',
    re.IGNORECASE
)
# Короткое тире (–) вместо длинного (—) между словами — ловит и двойные пробелы
EN_DASH_AS_EM   = re.compile(
    r'(?<=[а-яёА-ЯЁa-zA-Z])\s+–\s+(?=[а-яёА-ЯЁa-zA-Z])',
    re.IGNORECASE
)
RANGE_HYPHEN    = re.compile(r'\b(\d+)\s*-\s*(\d+)\b')
WRONG_QUOTES    = re.compile(r'"[^"]{1,80}"')
PERCENT_SPACE   = re.compile(r'\d\s+%')
NO_SPACE_UNIT   = re.compile(r'(\d)(млн|млрд|тыс)(?!\.)', re.IGNORECASE)
RUB_NO_DOT      = re.compile(r'\bруб\b(?!\.)', re.IGNORECASE)
MIN_NO_DOT      = re.compile(r'\bмин\b(?!\.)', re.IGNORECASE)
CURRENCY_BEFORE = re.compile(r'[₽](?=\d)')
MULTIPLY_X      = re.compile(r'\d\s*[xXхХ]\s*\d')
TITLE_DOT       = re.compile(r'\.\s*$')
DASH_NO_SPACE   = re.compile(r'\S—\S')
LARGE_NUMBER    = re.compile(r'(?<!\d)(\d{5,})(?!\d)')
YO_WORDS        = re.compile(
    r'\b(все|еще|ее|свое|своего|своих|своим|своими|своей|своему'
    r'|мое|твое|нее|него|нему|ней|поэтому|затем)\b',
    re.IGNORECASE | re.UNICODE
)


def _get_para_lines(para) -> list:
    """
    Возвращает список строк параграфа с учётом ручных переносов <a:br/>.
    python-pptx не включает <a:br/> в para.text, поэтому читаем XML напрямую.
    """
    from pptx.oxml.ns import qn
    lines, current = [], []
    for child in para._p:
        if child.tag == qn('a:r'):
            t = child.find(qn('a:t'))
            current.append(t.text or '' if t is not None else '')
        elif child.tag == qn('a:br'):
            lines.append(''.join(current))
            current = []
    lines.append(''.join(current))
    return [l for l in lines if l.strip()]


def check_typography(text, shape_name, is_title=False, text_frame=None):
    issues = []

    # Висячие предлоги — читаем строки с учётом ручных переносов <a:br/>
    if text_frame:
        seen = set()
        for para in text_frame.paragraphs:
            for line_text in _get_para_lines(para):
                lt = line_text.strip()
                if lt and HANGING_RE.search(lt):
                    word = lt.split()[-1]
                    key = lt[-30:]
                    if key not in seen:
                        seen.add(key)
                        issues.append({"code": "HANGING_PREPOSITION", "severity": "warning",
                            "message": f"Висячий предлог «{word}» — неразрывный пробел (Ctrl+Shift+Пробел)",
                            "line": lt[:80]})

    lines_list = text.split("\n")

    # ── Проверка стыков строк: «слово\n- следующее» или «слово\n– следующее» ──
    # Случай: ручной перенос стоит перед тире/дефисом
    _WORD_END   = re.compile(r'[а-яёА-ЯЁa-zA-Z\d]$')
    _DASH_NEXT  = re.compile(r'^-\s+[а-яёА-ЯЁa-zA-Z\d]')
    _EN_NEXT    = re.compile(r'^–\s+[а-яёА-ЯЁa-zA-Z\d]')

    for i in range(len(lines_list) - 1):
        curr = lines_list[i].rstrip()
        nxt  = lines_list[i + 1].lstrip()
        if not curr or not nxt:
            continue
        if _WORD_END.search(curr):
            preview = f"{curr[-35:]} / {nxt[:35]}"
            if _DASH_NEXT.match(nxt):
                issues.append({"code": "HYPHEN_AS_DASH", "severity": "error",
                    "message": "Дефис вместо тире — нужно «слово — слово»",
                    "line": preview[:80]})
            elif _EN_NEXT.match(nxt):
                issues.append({"code": "EN_DASH_AS_EM", "severity": "error",
                    "message": "Короткое тире (–) вместо длинного (—) между словами",
                    "line": preview[:80]})

    for line in lines_list:
        s = line.strip()
        if not s:
            continue
        if DOUBLE_SPACE.search(line):
            issues.append({"code": "DOUBLE_SPACE", "severity": "warning",
                "message": "Двойной пробел", "line": s[:80]})
        if HYPHEN_AS_DASH.search(line):
            issues.append({"code": "HYPHEN_AS_DASH", "severity": "error",
                "message": "Дефис вместо тире — нужно «слово — слово»", "line": s[:80]})
        if EN_DASH_AS_EM.search(line):
            issues.append({"code": "EN_DASH_AS_EM", "severity": "error",
                "message": "Короткое тире (–) вместо длинного (—) между словами",
                "line": s[:80]})
        for m in RANGE_HYPHEN.findall(line):
            issues.append({"code": "RANGE_HYPHEN", "severity": "error",
                "message": f"Диапазон «{m[0]}-{m[1]}» → «{m[0]}–{m[1]}» (короткое тире)", "line": s[:80]})
        if WRONG_QUOTES.search(line):
            issues.append({"code": "WRONG_QUOTES", "severity": "error",
                "message": 'Прямые кавычки " " → «ёлочки»', "line": s[:80]})
        if PERCENT_SPACE.search(line):
            issues.append({"code": "PERCENT_SPACE", "severity": "warning",
                "message": "«26%» без пробела, не «26 %»", "line": s[:80]})
        for m in NO_SPACE_UNIT.findall(line):
            issues.append({"code": "NO_SPACE_UNIT", "severity": "warning",
                "message": f"«{m[0]}{m[1]}» → «{m[0]} {m[1]}»", "line": s[:80]})
        if RUB_NO_DOT.search(line):
            issues.append({"code": "RUB_NO_DOT", "severity": "warning",
                "message": "«руб» → «руб.»", "line": s[:80]})
        if MIN_NO_DOT.search(line):
            issues.append({"code": "MIN_NO_DOT", "severity": "warning",
                "message": "«мин» → «мин.»", "line": s[:80]})
        if CURRENCY_BEFORE.search(line):
            issues.append({"code": "CURRENCY_BEFORE", "severity": "error",
                "message": "₽ ставится после числа: «1 000 ₽»", "line": s[:80]})
        if MULTIPLY_X.search(line):
            issues.append({"code": "MULTIPLY_X", "severity": "warning",
                "message": "Для умножения — знак «×», не буква X", "line": s[:80]})
        if is_title and TITLE_DOT.search(s):
            issues.append({"code": "TITLE_DOT", "severity": "error",
                "message": "Точка в конце заголовка не ставится", "line": s[:80]})
        if DASH_NO_SPACE.search(line):
            issues.append({"code": "DASH_NO_SPACE", "severity": "warning",
                "message": "Длинное тире (—) отбивается пробелами с двух сторон", "line": s[:80]})
        for num in LARGE_NUMBER.findall(line):
            if re.match(r'^(19|20)\d\d$', num):
                continue
            issues.append({"code": "LARGE_NUMBER", "severity": "warning",
                "message": f"«{num}» → «{format_number(num)}»", "line": s[:80]})
        for m in YO_WORDS.findall(s):
            if 'ё' not in m.lower():
                issues.append({"code": "MISSING_YO", "severity": "warning",
                    "message": f"Возможно пропущена «ё»: «{m}»", "line": s[:80]})
    return issues


# ─────────────────────────────────────────────
# БЛОК 2: ШРИФТЫ И ТЕКСТ
# ─────────────────────────────────────────────

def check_title_length(shape):
    """1.2 Длина заголовка"""
    if not shape.has_text_frame:
        return []
    text = shape.text_frame.text.strip()
    if len(text) > MAX_TITLE_LEN:
        return [{"code": "TITLE_TOO_LONG", "severity": "warning",
            "message": f"Заголовок слишком длинный: {len(text)} символов (рекомендуется до {MAX_TITLE_LEN})",
            "line": text[:80]}]
    return []


def check_font_sizes_on_slide(slide) -> list:
    """
    1.3 Шрифты: слишком маленький, слишком много разных размеров.
    Возвращает список ошибок для слайда.
    """
    issues = []
    all_sizes = []

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        is_title = is_title_shape(shape)
        sizes = get_font_sizes(shape)
        all_sizes.extend(sizes)

        for sz in sizes:
            threshold = MIN_TITLE_FONT_PT if is_title else MIN_BODY_FONT_PT
            if sz < threshold and sz > 0:
                issues.append({"code": "FONT_TOO_SMALL", "severity": "warning",
                    "message": f"«{shape.name}»: шрифт {sz}pt — слишком мелко (мин. {threshold}pt)"})

    # Слишком много разных размеров (фильтруем декоративные >100pt)
    filtered_sizes = [s for s in all_sizes if s <= 100]
    unique_sizes = set(round(s / 2) * 2 for s in filtered_sizes)
    if len(unique_sizes) > MAX_FONT_SIZES_SLIDE:
        issues.append({"code": "TOO_MANY_FONT_SIZES", "severity": "warning",
            "message": f"Много разных размеров шрифта: {sorted(set(round(s) for s in filtered_sizes))} pt"})

    return issues



# Таблица рекомендуемого межстрочного по гайду Яндекс Рекламы
# (min_pt, max_pt, рекомендуемый_ls, допуск)
LS_GUIDE = [
    (140, 9999, 0.80, 0.04),
    (115,  140, 0.82, 0.04),
    ( 90,  115, 0.82, 0.04),
    ( 65,   90, 0.84, 0.05),
    ( 50,   65, 0.89, 0.06),
    ( 38,   50, 1.00, 0.06),
    ( 25,   38, 1.04, 0.07),
    (  0,   25, 1.10, 0.07),
]

def _recommended_ls(font_pt: float):
    for min_pt, max_pt, rec, tol in LS_GUIDE:
        if min_pt <= font_pt < max_pt:
            return rec, tol
    return None


def check_line_spacing(slide) -> list:
    """1.5 Межстрочный интервал — базовые ограничения + проверка по гайду"""
    issues = []
    seen_guide = set()

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            ls = para.line_spacing
            if ls is None:
                continue
            try:
                ls_val = float(ls)
            except Exception:
                continue
            if ls_val > 10:   # EMU/Pt значение — пропускаем
                continue

            # Базовые ограничения
            if ls_val < LINE_SPACING_MIN:
                issues.append({"code": "LINE_SPACING_TIGHT", "severity": "warning",
                    "message": f"«{shape.name}»: межстрочный {ls_val:.2f} — слишком плотно"})
                continue
            if ls_val > LINE_SPACING_MAX:
                issues.append({"code": "LINE_SPACING_LOOSE", "severity": "warning",
                    "message": f"«{shape.name}»: межстрочный {ls_val:.2f} — слишком широко"})
                continue

            # Проверка по таблице гайда
            sizes = [r.font.size for r in para.runs if r.font.size]
            if not sizes:
                continue
            font_pt = max(sizes) / 12700
            if font_pt > 100:   # декоративные шрифты пропускаем
                continue
            rec = _recommended_ls(font_pt)
            if not rec:
                continue
            rec_ls, tol = rec
            if abs(ls_val - rec_ls) > tol:
                key = (shape.name, round(font_pt), round(ls_val, 2))
                if key not in seen_guide:
                    seen_guide.add(key)
                    issues.append({"code": "LINE_SPACING_GUIDE", "severity": "warning",
                        "message": (
                            f"«{shape.name}»: интерлиньяж {ls_val:.2f} "
                            f"при {font_pt:.0f}pt (по гайду: {rec_ls})"
                        )})

    return issues





# ─────────────────────────────────────────────
# ШРИФТЫ: СЕМЕЙСТВО, ДОПУСТИМЫЕ РАЗМЕРЫ
# ─────────────────────────────────────────────

# Разрешённые размеры шрифта (pt)
ALLOWED_FONT_SIZES_PT = {20, 24, 28, 32, 36, 40, 48, 56, 72, 96, 120, 148}
FONT_SIZE_SNAP_TOL    = 2   # ±2pt — считается допустимым

# Максимум разных шрифтовых семейств на слайде
MAX_FONT_FAMILIES = 2


def _normalize_font_name(name: str | None) -> str | None:
    if not name:
        return None
    name = re.sub(r"\s+", " ", str(name)).strip()
    return name or None


def load_theme_fonts(pptx_path: str) -> dict:
    """Читает шрифты темы из ppt/theme/theme*.xml."""
    theme_fonts = {
        "major_latin": None,
        "minor_latin": None,
        "major_ea": None,
        "minor_ea": None,
        "major_cs": None,
        "minor_cs": None,
    }
    try:
        with zipfile.ZipFile(pptx_path) as zf:
            theme_files = sorted(
                n for n in zf.namelist()
                if n.startswith("ppt/theme/theme") and n.endswith(".xml")
            )
            if not theme_files:
                return theme_fonts

            root = etree.fromstring(zf.read(theme_files[0]))
            ns = {"a": NS}
            for scheme, prefix in (("majorFont", "major"), ("minorFont", "minor")):
                for script, suffix in (("latin", "latin"), ("ea", "ea"), ("cs", "cs")):
                    val = root.xpath(
                        f"string(.//a:themeElements/a:fontScheme/a:{scheme}/a:{script}/@typeface)",
                        namespaces=ns,
                    )
                    theme_fonts[f"{prefix}_{suffix}"] = _normalize_font_name(val)
    except Exception:
        pass
    return theme_fonts


def _first_xml_typeface(el) -> str | None:
    if el is None:
        return None
    for tag in ("latin", "ea", "cs", "sym"):
        child = el.find(f"{{{NS}}}{tag}")
        if child is not None:
            val = _normalize_font_name(child.get("typeface"))
            if val:
                return val
    return None


def _font_name_from_run(run) -> str | None:
    direct = _normalize_font_name(run.font.name)
    if direct:
        return direct
    return _first_xml_typeface(getattr(run._r, 'rPr', None))


def _theme_fallback_font(is_title: bool, theme_fonts: dict | None) -> str | None:
    theme_fonts = theme_fonts or {}
    if is_title:
        return (
            theme_fonts.get("major_latin")
            or theme_fonts.get("major_ea")
            or theme_fonts.get("major_cs")
        )
    return (
        theme_fonts.get("minor_latin")
        or theme_fonts.get("minor_ea")
        or theme_fonts.get("minor_cs")
    )


def get_shape_font_families(shape, theme_fonts: dict | None = None, is_title: bool | None = None) -> list:
    families = []
    if not shape.has_text_frame:
        return families
    if is_title is None:
        is_title = is_title_shape(shape)

    for para in shape.text_frame.paragraphs:
        para_fallback = _first_xml_typeface(para._p.find(f"{{{NS}}}pPr"))
        for run in para.runs:
            if not run.text.strip():
                continue
            family = _font_name_from_run(run) or para_fallback or _theme_fallback_font(is_title, theme_fonts)
            family = _normalize_font_name(family)
            if family:
                families.append(family)

    if not families:
        fallback = _theme_fallback_font(is_title, theme_fonts)
        if fallback and shape.text_frame.text.strip():
            families.append(fallback)

    return families


def check_font_size_allowed(slide) -> list:
    """Флагит размеры шрифта не из разрешённого списка."""
    issues = []
    seen = set()
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not run.font.size:
                    continue
                pt = round(run.font.size / 12700)
                if pt in seen or pt == 0 or pt > 160:
                    continue
                seen.add(pt)
                if not any(abs(pt - a) <= FONT_SIZE_SNAP_TOL for a in ALLOWED_FONT_SIZES_PT):
                    closest = min(ALLOWED_FONT_SIZES_PT, key=lambda a: abs(a - pt))
                    issues.append({
                        "code": "FONT_SIZE_NOT_ALLOWED", "severity": "warning",
                        "message": (
                            f"Нестандартный размер шрифта {pt}pt "
                            f"(ближайший допустимый: {closest}pt)"
                        ),
                    })
    return issues


def check_font_families(slide, theme_fonts: dict | None = None) -> list:
    """Флагит смешение более MAX_FONT_FAMILIES семейств шрифтов на слайде."""
    issues = []
    families = set()
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        families.update(get_shape_font_families(shape, theme_fonts=theme_fonts))
    if len(families) > MAX_FONT_FAMILIES:
        issues.append({
            "code": "FONT_MIX", "severity": "warning",
            "message": f"Смешение шрифтов: {', '.join(sorted(families))}",
        })
    return issues


# ─────────────────────────────────────────────
# ЗАГОЛОВКИ: ВЫРАВНИВАНИЕ И НАЧЕРТАНИЕ
# ─────────────────────────────────────────────

_ALIGN_NAMES = {1: "по левому краю", 2: "по центру",
                3: "по правому краю", 4: "по ширине"}


def _get_title_align(slide):
    """Возвращает выравнивание первого абзаца заголовка или None."""
    for shape in slide.shapes:
        if is_title_shape(shape) and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if para.alignment is not None:
                    return int(para.alignment)
            return 1  # left by default
    return None


def _get_title_bold(slide):
    """Возвращает bold первого рана заголовка или None."""
    for shape in slide.shapes:
        if is_title_shape(shape) and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.bold is not None:
                        return run.font.bold
    return None


def check_title_alignment_consistency(slides) -> dict:
    """Проверяет что выравнивание заголовков одинаково на всех слайдах."""
    results = defaultdict(list)
    aligns = {i + 1: _get_title_align(s) for i, s in enumerate(slides)}
    aligns = {sn: a for sn, a in aligns.items() if a is not None}
    if len(aligns) < 2:
        return {}
    mode = Counter(aligns.values()).most_common(1)[0][0]
    for sn, align in aligns.items():
        if align != mode:
            results[sn].append({
                "code": "TITLE_ALIGN_INCONSISTENT", "severity": "warning",
                "message": (
                    f"Заголовок: {_ALIGN_NAMES.get(align, align)} "
                    f"(большинство — {_ALIGN_NAMES.get(mode, mode)})"
                ),
            })
    return dict(results)


def check_title_bold_consistency(slides) -> dict:
    """Проверяет что начертание (bold/regular) заголовков одинаково."""
    results = defaultdict(list)
    bolds = {i + 1: _get_title_bold(s) for i, s in enumerate(slides)}
    bolds = {sn: b for sn, b in bolds.items() if b is not None}
    if len(bolds) < 2:
        return {}
    mode = Counter(bolds.values()).most_common(1)[0][0]
    for sn, bold in bolds.items():
        if bold != mode:
            results[sn].append({
                "code": "TITLE_BOLD_INCONSISTENT", "severity": "warning",
                "message": (
                    f"Начертание заголовка: {'жирный' if bold else 'обычный'} "
                    f"(большинство — {'жирный' if mode else 'обычный'})"
                ),
            })
    return dict(results)


# ─────────────────────────────────────────────
# ПОЛЯ: ОПРЕДЕЛЯЕМ ПО СЛАЙДУ 2
# ─────────────────────────────────────────────

CONTENT_MARGIN_TOL = int(0.5 / 2.54 * 914400)  # 0.5 см допуск


def _infer_left_margin(slide) -> int | None:
    """Минимальный left текстовых элементов на слайде (не считая крайне левых)."""
    min_left = None
    for shape in slide.shapes:
        if not shape.has_text_frame or not shape.text_frame.text.strip():
            continue
        if shape.left is None or emu_to_cm(shape.left) < 0.5:
            continue  # декоративные у левого края пропускаем
        if min_left is None or shape.left < min_left:
            min_left = shape.left
    return min_left


def check_content_margins(slides) -> dict:
    """
    Использует слайд 2 как эталон левого поля.
    Флагит текстовые элементы на слайдах 3+ которые левее эталона.
    """
    results = defaultdict(list)
    if len(slides) < 2:
        return {}

    ref_left = _infer_left_margin(slides[1])  # слайд 2 (индекс 1)
    if ref_left is None:
        return {}

    for slide_idx in range(2, len(slides)):  # начиная со слайда 3
        slide_num = slide_idx + 1
        for shape in slides[slide_idx].shapes:
            if not shape.has_text_frame or not shape.text_frame.text.strip():
                continue
            if shape.left is None:
                continue
            if shape.left < ref_left - CONTENT_MARGIN_TOL:
                diff = emu_to_cm(ref_left - shape.left)
                results[slide_num].append({
                    "code": "CONTENT_MARGIN_LEFT", "severity": "warning",
                    "message": (
                        f"«{shape.name}» левее поля на {diff:.1f} см "
                        f"(поле определено по слайду 2)"
                    ),
                })
    return dict(results)


def check_bullet_consistency(slide) -> list:
    """1.6 Буллиты: одинаковый тип и отступ внутри блоков"""
    issues = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        paras = [p for p in shape.text_frame.paragraphs if p.text.strip()]
        if len(paras) < 2:
            continue

        bullet_types = [get_bullet_type(p) for p in paras]
        levels       = [p.level for p in paras]

        # Разные типы буллитов внутри одного блока
        non_none = [t for t in bullet_types if t != "none" and t != "inherit"]
        if len(set(non_none)) > 1:
            issues.append({"code": "BULLET_INCONSISTENT", "severity": "warning",
                "message": f"«{shape.name}»: разные типы буллитов в одном блоке"})

        # Пропущенные уровни отступа (1 → 3 без 2)
        used_levels = sorted(set(levels))
        if len(used_levels) > 1:
            for i in range(len(used_levels) - 1):
                if used_levels[i+1] - used_levels[i] > 1:
                    issues.append({"code": "BULLET_LEVEL_GAP", "severity": "warning",
                        "message": f"«{shape.name}»: пропущен уровень отступа ({used_levels[i]} → {used_levels[i+1]})"})
    return issues


def extract_title_font_meta(slide, theme_fonts: dict | None = None) -> dict | None:
    """Извлекает размер и семейство шрифта заголовка для межслайдовой проверки."""
    for shape in slide.shapes:
        if is_title_shape(shape) and shape.has_text_frame:
            sizes = get_font_sizes(shape)
            families = get_shape_font_families(shape, theme_fonts=theme_fonts, is_title=True)
            family_counter = Counter(families)
            return {
                "font_sizes": sizes,
                "max_size": max(sizes) if sizes else None,
                "families": sorted(set(families)),
                "main_family": family_counter.most_common(1)[0][0] if family_counter else None,
            }
    return None


def extract_body_font_meta(slide, theme_fonts: dict | None = None) -> dict | None:
    """Извлекает основной шрифт тела слайда (не заголовок)."""
    all_sizes = []
    families = []
    for shape in slide.shapes:
        if is_title_shape(shape) or not shape.has_text_frame:
            continue
        sizes = get_font_sizes(shape)
        all_sizes.extend(sizes)
        families.extend(get_shape_font_families(shape, theme_fonts=theme_fonts, is_title=False))
    if not all_sizes and not families:
        return None
    size_counter = Counter(round(s) for s in all_sizes)
    family_counter = Counter(families)
    main_size = size_counter.most_common(1)[0][0] if size_counter else None
    main_family = family_counter.most_common(1)[0][0] if family_counter else None
    return {
        "main_size": main_size,
        "families": sorted(set(families)),
        "main_family": main_family,
    }


def _mode_value(metas: list[dict | None], key: str):
    values = [m.get(key) for m in metas if m and m.get(key)]
    if not values:
        return None
    return Counter(values).most_common(1)[0][0]


def check_font_consistency_across_slides(all_title_metas, all_body_metas) -> dict:
    """
    1.4 Консистентность шрифтов между слайдами:
    заголовки и основной текст должны иметь одинаковый размер и семейство.
    """
    results = defaultdict(list)

    ref_title_size = _mode_value(all_title_metas, "max_size")
    ref_body_size = _mode_value(all_body_metas, "main_size")
    ref_title_family = _mode_value(all_title_metas, "main_family")
    ref_body_family = _mode_value(all_body_metas, "main_family")

    for idx, meta in enumerate(all_title_metas):
        slide_num = idx + 1
        if not meta:
            continue
        if ref_title_size and meta.get("max_size"):
            if abs(meta["max_size"] - ref_title_size) > FONT_SIZE_TOLERANCE:
                results[slide_num].append({
                    "code": "TITLE_FONT_INCONSISTENT", "severity": "error",
                    "message": f"Размер заголовка {meta['max_size']}pt vs {ref_title_size}pt (большинство слайдов)"
                })
        if ref_title_family and meta.get("main_family") and meta["main_family"] != ref_title_family:
            results[slide_num].append({
                "code": "TITLE_FONT_FAMILY_INCONSISTENT", "severity": "error",
                "message": f"Шрифт заголовка «{meta['main_family']}» vs «{ref_title_family}» (большинство слайдов)"
            })

    for idx, meta in enumerate(all_body_metas):
        slide_num = idx + 1
        if not meta:
            continue
        if ref_body_size and meta.get("main_size"):
            if abs(meta["main_size"] - ref_body_size) > FONT_SIZE_TOLERANCE:
                results[slide_num].append({
                    "code": "BODY_FONT_INCONSISTENT", "severity": "warning",
                    "message": f"Основной шрифт текста {meta['main_size']}pt vs {ref_body_size}pt (большинство слайдов)"
                })
        if ref_body_family and meta.get("main_family") and meta["main_family"] != ref_body_family:
            results[slide_num].append({
                "code": "BODY_FONT_FAMILY_INCONSISTENT", "severity": "warning",
                "message": f"Шрифт основного текста «{meta['main_family']}» vs «{ref_body_family}» (большинство слайдов)"
            })

    return dict(results)


# ─────────────────────────────────────────────
# БЛОК 3: РАСПОЛОЖЕНИЕ ЭЛЕМЕНТОВ
# ─────────────────────────────────────────────

def check_element_overlap(slide) -> list:
    """2.4 Перекрытие элементов (bounding box)"""
    issues = []
    threshold = int(OVERLAP_MIN_CM / 2.54 * 914400)

    # Пропускаем фигуры-декор (линии, стрелки, скобки) и очень маленькие фигуры
    SKIP_OVERLAP_KEYWORDS = ("линия", "скобка", "стрелка", "соединит", "полилиния")

    shapes = [
        (s.name, s.left, s.top, s.left + s.width, s.top + s.height)
        for s in slide.shapes
        if s.left >= 0 and s.top >= 0 and s.width > 0 and s.height > 0
        and not s.name.startswith("think-cell")
        and not any(kw in s.name.lower() for kw in SKIP_OVERLAP_KEYWORDS)
        and s.width > int(1 / 2.54 * 914400)   # шире 1 см
        and s.height > int(0.5 / 2.54 * 914400) # выше 0.5 см
    ]

    seen_pairs = set()
    for i, (n1, l1, t1, r1, b1) in enumerate(shapes):
        for n2, l2, t2, r2, b2 in shapes[i+1:]:
            if l1 < r2 and r1 > l2 and t1 < b2 and b1 > t2:
                ow = min(r1, r2) - max(l1, l2)
                oh = min(b1, b2) - max(t1, t2)
                if ow > threshold and oh > threshold:
                    pair = tuple(sorted([n1, n2]))
                    if pair not in seen_pairs:
                        seen_pairs.add(pair)
                        issues.append({"code": "ELEMENT_OVERLAP", "severity": "error",
                            "message": f"Перекрытие: «{n1}» и «{n2}» ({emu_to_cm(ow)}×{emu_to_cm(oh)} см)"})
    return issues


def check_near_alignment(slide) -> list:
    """
    2.1 «Почти выровнено»: два элемента с разницей left/top
    менее SNAP_THRESHOLD_CM, но больше POSITION_TOLERANCE.
    """
    issues = []
    tol_lo = POSITION_TOLERANCE
    tol_hi = int(SNAP_THRESHOLD_CM / 2.54 * 914400)

    shapes = [
        (s.name, s.left, s.top)
        for s in slide.shapes
        if not s.name.startswith("think-cell") and not s.name.startswith("Google Shape")
    ]

    seen_pairs = set()
    for i, (n1, l1, t1) in enumerate(shapes):
        for n2, l2, t2 in shapes[i+1:]:
            # Горизонталь
            dl = abs(l1 - l2)
            if tol_lo < dl < tol_hi:
                pair = tuple(sorted([n1, n2]) + ["left"])
                if pair not in seen_pairs:
                    seen_pairs.add(pair)
                    issues.append({"code": "NEAR_ALIGNED", "severity": "warning",
                        "message": f"«{n1}» и «{n2}» почти выровнены по горизонтали (разница {emu_to_cm(dl)} см)"})
            # Вертикаль
            dt = abs(t1 - t2)
            if tol_lo < dt < tol_hi:
                pair = tuple(sorted([n1, n2]) + ["top"])
                if pair not in seen_pairs:
                    seen_pairs.add(pair)
                    issues.append({"code": "NEAR_ALIGNED", "severity": "warning",
                        "message": f"«{n1}» и «{n2}» почти выровнены по вертикали (разница {emu_to_cm(dt)} см)"})
    return issues


def check_element_spacing(slide) -> list:
    """2.3 Слишком маленькое расстояние между элементами (но не перекрытие)"""
    issues = []
    MIN_GAP = int(0.2 / 2.54 * 914400)   # 0.2 cm минимальный зазор

    shapes = [
        (s.name, s.left, s.top, s.left + s.width, s.top + s.height)
        for s in slide.shapes
        if s.left >= 0 and s.top >= 0 and s.width > 0 and s.height > 0
        and not s.name.startswith("think-cell")
        and not s.name.startswith("Google Shape")
    ]

    seen = set()
    for i, (n1, l1, t1, r1, b1) in enumerate(shapes):
        for n2, l2, t2, r2, b2 in shapes[i+1:]:
            # Только если они рядом (не перекрываются, но близко)
            h_gap = max(l2 - r1, l1 - r2, 0)
            v_gap = max(t2 - b1, t1 - b2, 0)
            if h_gap == 0 and 0 < v_gap < MIN_GAP:
                pair = tuple(sorted([n1, n2]))
                if pair not in seen:
                    seen.add(pair)
                    issues.append({"code": "TIGHT_SPACING", "severity": "warning",
                        "message": f"«{n1}» и «{n2}» стоят очень близко по вертикали ({emu_to_cm(v_gap)} см)"})
            if v_gap == 0 and 0 < h_gap < MIN_GAP:
                pair = tuple(sorted([n1, n2]))
                if pair not in seen:
                    seen.add(pair)
                    issues.append({"code": "TIGHT_SPACING", "severity": "warning",
                        "message": f"«{n1}» и «{n2}» стоят очень близко по горизонтали ({emu_to_cm(h_gap)} см)"})
    return issues


# ─────────────────────────────────────────────
# БЛОК 4: ИЗОБРАЖЕНИЯ
# ─────────────────────────────────────────────

ASPECT_RATIO_TOLERANCE = 0.15   # допустимое отклонение соотношения сторон

def check_image_aspect_ratio(slide) -> list:
    """3.2 Растягивание изображений"""
    issues = []
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        try:
            # Получаем оригинальные размеры из blob
            from PIL import Image
            import io
            img_blob = shape.image.blob
            img = Image.open(io.BytesIO(img_blob))
            orig_w, orig_h = img.size
            orig_ratio = orig_w / orig_h
            curr_ratio = shape.width / shape.height
            if abs(curr_ratio - orig_ratio) / orig_ratio > ASPECT_RATIO_TOLERANCE:
                issues.append({"code": "IMAGE_DISTORTED", "severity": "warning",
                    "message": (
                        f"«{shape.name}»: изображение искажено — "
                        f"текущее соотношение {curr_ratio:.2f}, оригинал {orig_ratio:.2f}"
                    )})
        except Exception:
            pass
    return issues


def extract_image_metas(slide) -> list:
    return [
        {"name": s.name, "left": s.left, "top": s.top,
         "width": s.width, "height": s.height}
        for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]


def check_image_consistency(all_image_metas: list) -> dict:
    """3.1 и 3.3 Одинаковые изображения должны совпадать по размеру и позиции"""
    errors = defaultdict(list)
    single = [(i + 1, m[0]) for i, m in enumerate(all_image_metas) if len(m) == 1]
    if len(single) < 2:
        return {}
    _, ref = single[0]
    for slide_num, img in single[1:]:
        if abs(img["width"] - ref["width"]) > SIZE_TOLERANCE:
            errors[slide_num].append({"code": "IMAGE_SIZE_MISMATCH", "severity": "warning",
                "message": f"Ширина изображения {emu_to_cm(img['width'])} см vs {emu_to_cm(ref['width'])} см (слайд 1)"})
        if abs(img["height"] - ref["height"]) > SIZE_TOLERANCE:
            errors[slide_num].append({"code": "IMAGE_SIZE_MISMATCH", "severity": "warning",
                "message": f"Высота изображения {emu_to_cm(img['height'])} см vs {emu_to_cm(ref['height'])} см (слайд 1)"})
        if abs(img["left"] - ref["left"]) > POSITION_TOLERANCE:
            errors[slide_num].append({"code": "IMAGE_POS_MISMATCH", "severity": "warning",
                "message": f"Изображение сдвинуто горизонтально: {emu_to_cm(img['left'])} vs {emu_to_cm(ref['left'])} см"})
        if abs(img["top"] - ref["top"]) > POSITION_TOLERANCE:
            errors[slide_num].append({"code": "IMAGE_POS_MISMATCH", "severity": "warning",
                "message": f"Изображение сдвинуто вертикально: {emu_to_cm(img['top'])} vs {emu_to_cm(ref['top'])} см"})
    return dict(errors)


# ─────────────────────────────────────────────
# БЛОК 5: ЦВЕТА
# ─────────────────────────────────────────────

def check_color_count(slide) -> list:
    """5.1 Слишком много уникальных цветов на слайде"""
    all_colors = set()
    for shape in slide.shapes:
        all_colors.update(get_shape_colors(shape))
    # Убираем очень светлые (белый, почти белый) и очень тёмные (чёрный)
    filtered = {c for c in all_colors if c not in ("000000", "FFFFFF", "FEFEFE", "FDFDFD")}
    if len(filtered) > MAX_COLORS_SLIDE:
        return [{"code": "TOO_MANY_COLORS", "severity": "warning",
            "message": f"Много цветов на слайде: {len(filtered)} уникальных (рекомендуется до {MAX_COLORS_SLIDE})"}]
    return []


# ─────────────────────────────────────────────
# БЛОК 6: ТАБЛИЦЫ
# ─────────────────────────────────────────────

def check_tables(slide) -> list:
    """6.x Проверка таблиц"""
    issues = []
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        table = shape.table
        nrows = len(table.rows)
        ncols = len(table.columns)

        # 6.3 Перегруженность
        if nrows > MAX_TABLE_ROWS:
            issues.append({"code": "TABLE_TOO_MANY_ROWS", "severity": "warning",
                "message": f"Таблица «{shape.name}»: {nrows} строк — много для презентации (рекомендуется до {MAX_TABLE_ROWS})"})
        if ncols > MAX_TABLE_COLS:
            issues.append({"code": "TABLE_TOO_MANY_COLS", "severity": "warning",
                "message": f"Таблица «{shape.name}»: {ncols} столбцов (рекомендуется до {MAX_TABLE_COLS})"})

        # 6.2 Повторяющиеся заголовки
        if nrows > 0:
            header_cells = [c.text.strip() for c in table.rows[0].cells]
            if len(set(header_cells)) < len([h for h in header_cells if h]):
                issues.append({"code": "TABLE_DUPLICATE_HEADERS", "severity": "warning",
                    "message": f"Таблица «{shape.name}»: повторяющиеся заголовки столбцов"})

        # 6.1 Выравнивание: числа должны быть справа
        for row in table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    text = para.text.strip()
                    if re.match(r'^\d', text) and para.alignment == PP_ALIGN.LEFT:
                        issues.append({"code": "TABLE_NUMBER_ALIGN", "severity": "warning",
                            "message": f"Таблица «{shape.name}»: числа лучше выравнивать по правому краю"})
                        break

    return issues


# ─────────────────────────────────────────────
# БЛОК 7: КОНСИСТЕНТНОСТЬ ЭЛЕМЕНТОВ МЕЖДУ СЛАЙДАМИ
# ─────────────────────────────────────────────


# ─────────────────────────────────────────────
# ПРОВЕРКА СКРУГЛЕНИЙ ФИГУР
# ─────────────────────────────────────────────

_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

def get_roundness(shape):
    """Возвращает скругление угла фигуры (0–50000) или None."""
    import re as _re
    try:
        avLst = shape._element.find(f".//{{{_NS_A}}}avLst")
        if avLst is None:
            return None
        for gd in avLst:
            if gd.get("name") == "adj":
                m = _re.search(r"\d+", gd.get("fmla", ""))
                if m:
                    return int(m.group())
    except Exception:
        pass
    return None


def check_roundness_consistency(slides) -> dict:
    """Сравнивает скругления одноимённых фигур между слайдами."""
    results = defaultdict(list)
    reference = {}
    ref_slide  = {}

    for slide_idx, slide in enumerate(slides):
        for shape in slide.shapes:
            name = shape.name
            if not name or name.startswith("think-cell") or name.startswith("Google Shape"):
                continue
            r = get_roundness(shape)
            if r is None:
                continue
            if name not in reference:
                reference[name] = r
                ref_slide[name] = slide_idx + 1

    for slide_idx, slide in enumerate(slides):
        slide_num = slide_idx + 1
        seen = set()
        for shape in slide.shapes:
            name = shape.name
            if name in seen:
                continue
            r = get_roundness(shape)
            if r is None or name not in reference or ref_slide[name] == slide_num:
                continue
            ref_r   = reference[name]
            ref_num = ref_slide[name]
            if r != ref_r:
                seen.add(name)
                results[slide_num].append({
                    "code": "ROUNDNESS_MISMATCH", "severity": "warning",
                    "message": (
                        f"«{name}»: скругление {round(r/500)}% "
                        f"vs {round(ref_r/500)}% (слайд {ref_num})"
                    ),
                })

    return dict(results)


def extract_slide_shapes_meta(slide) -> dict:
    meta = {}
    for shape in slide.shapes:
        name = shape.name
        if not name or name.startswith("think-cell") or name.startswith("Google Shape"):
            continue
        key = make_shape_key(name, shape.left, shape.top)
        meta[key] = {
            "display_name": name,
            "left":       shape.left,
            "top":        shape.top,
            "width":      shape.width,
            "height":     shape.height,
            "font_sizes": get_font_sizes(shape),
            "is_title":   is_title_shape(shape),
        }
    return meta


def check_all_shapes_consistency(all_slide_metas: list) -> dict:
    """Сравнивает позиции и шрифты ТОЛЬКО заголовков между слайдами."""
    results = defaultdict(list)
    reference = {}
    ref_slide  = {}

    for slide_idx, meta in enumerate(all_slide_metas):
        for key, props in meta.items():
            if not props.get("is_title"):
                continue
            if key not in reference:
                reference[key] = props
                ref_slide[key]  = slide_idx + 1

    for slide_idx, meta in enumerate(all_slide_metas):
        slide_num = slide_idx + 1
        for key, props in meta.items():
            if not props.get("is_title"):
                continue
            ref = reference.get(key)
            if not ref or ref_slide[key] == slide_num:
                continue
            name    = props["display_name"]
            ref_num = ref_slide[key]

            if abs(props["top"] - ref["top"]) > POSITION_TOLERANCE:
                diff = emu_to_cm(abs(props["top"] - ref["top"]))
                results[slide_num].append({"code": "SHAPE_TOP_MISMATCH", "severity": "error",
                    "message": f"Заголовок «{name}» сдвинут на {diff} см по вертикали (эталон — слайд {ref_num})"})

            if abs(props["left"] - ref["left"]) > POSITION_TOLERANCE:
                diff = emu_to_cm(abs(props["left"] - ref["left"]))
                results[slide_num].append({"code": "SHAPE_LEFT_MISMATCH", "severity": "error",
                    "message": f"Заголовок «{name}» сдвинут на {diff} см по горизонтали (эталон — слайд {ref_num})"})

            if props["font_sizes"] and ref["font_sizes"]:
                cur_max = max(props["font_sizes"])
                ref_max = max(ref["font_sizes"])
                if abs(cur_max - ref_max) > FONT_SIZE_TOLERANCE:
                    results[slide_num].append({"code": "SHAPE_FONT_MISMATCH", "severity": "error",
                        "message": f"Заголовок «{name}» шрифт {cur_max}pt вместо {ref_max}pt (слайд {ref_num})"})

    return dict(results)


# ─────────────────────────────────────────────
# БЛОК 8: ШАБЛОН
# ─────────────────────────────────────────────

def extract_layout_snapshot(filepath: str) -> list:
    prs = Presentation(filepath)
    return [
        [{"name": s.name, "left": s.left, "top": s.top, "width": s.width, "height": s.height}
         for s in slide.shapes]
        for slide in prs.slides
    ]


def compare_with_template(template_path: str, presentation_path: str) -> dict:
    template_snap = extract_layout_snapshot(template_path)
    prs = Presentation(presentation_path)
    results = defaultdict(list)
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        if slide_idx >= len(template_snap):
            continue
        tmpl_by_key = {make_shape_key(e["name"], e["left"], e["top"]): e
                       for e in template_snap[slide_idx]}
        for shape in slide.shapes:
            key = make_shape_key(shape.name, shape.left, shape.top)
            tmpl = tmpl_by_key.get(key)
            if not tmpl:
                continue
            if abs(shape.top - tmpl["top"]) > POSITION_TOLERANCE:
                results[slide_num].append({"code": "TMPL_TOP", "severity": "error",
                    "message": f"«{shape.name}» смещён вертикально: {emu_to_cm(shape.top)} см, шаблон {emu_to_cm(tmpl['top'])} см"})
            if abs(shape.left - tmpl["left"]) > POSITION_TOLERANCE:
                results[slide_num].append({"code": "TMPL_LEFT", "severity": "error",
                    "message": f"«{shape.name}» смещён горизонтально: {emu_to_cm(shape.left)} см, шаблон {emu_to_cm(tmpl['left'])} см"})
        tmpl_count  = len(template_snap[slide_idx])
        slide_count = len(list(slide.shapes))
        if abs(slide_count - tmpl_count) > 2:
            results[slide_num].append({"code": "TMPL_COUNT", "severity": "warning",
                "message": f"Элементов: {slide_count}, в шаблоне: {tmpl_count}"})
    return dict(results)


# ─────────────────────────────────────────────
# ГЛАВНАЯ ФУНКЦИЯ
# ─────────────────────────────────────────────

def analyze_presentation(filepath: str, template_path: str = None) -> dict:
    theme_fonts = load_theme_fonts(filepath)
    prs = Presentation(filepath)
    total_slides = len(prs.slides)
    slide_reports = {}
    all_slide_metas  = []
    all_image_metas  = []
    all_title_metas  = []
    all_body_metas   = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        issues = []
        slide_title = f"Слайд {slide_num}"

        # ── Типографика ──
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = get_shape_text(shape).strip()
            if not text:
                continue
            is_title = is_title_shape(shape)
            if is_title:
                t = shape.text_frame.text.strip()
                if t:
                    slide_title = t
                issues.extend(check_title_length(shape))
            issues.extend(check_typography(text, shape.name, is_title=is_title,
                                           text_frame=shape.text_frame))

        # ── Шрифты на слайде ──
        issues.extend(check_font_sizes_on_slide(slide))
        issues.extend(check_font_size_allowed(slide))
        issues.extend(check_font_families(slide, theme_fonts=theme_fonts))
        issues.extend(check_line_spacing(slide))
        issues.extend(check_bullet_consistency(slide))

        # ── Расположение ──


        # ── Изображения ──
        issues.extend(check_image_aspect_ratio(slide))

        # ── Цвета ──
        issues.extend(check_color_count(slide))

        # ── Таблицы ──
        issues.extend(check_tables(slide))

        slide_reports[slide_num] = {"title": slide_title, "issues": issues}

        # Метаданные для межслайдового анализа
        all_slide_metas.append(extract_slide_shapes_meta(slide))
        all_image_metas.append(extract_image_metas(slide))
        all_title_metas.append(extract_title_font_meta(slide, theme_fonts=theme_fonts))
        all_body_metas.append(extract_body_font_meta(slide, theme_fonts=theme_fonts))

    # ── Межслайдовые проверки ──
    for sn, issues in check_all_shapes_consistency(all_slide_metas).items():
        slide_reports[sn]["issues"].extend(issues)

    for sn, issues in check_image_consistency(all_image_metas).items():
        if sn in slide_reports:
            slide_reports[sn]["issues"].extend(issues)

    for sn, issues in check_font_consistency_across_slides(all_title_metas, all_body_metas).items():
        if sn in slide_reports:
            slide_reports[sn]["issues"].extend(issues)

    # Скругления фигур
    for sn, issues in check_roundness_consistency(list(prs.slides)).items():
        if sn in slide_reports:
            slide_reports[sn]["issues"].extend(issues)

    # Выравнивание и начертание заголовков
    for sn, issues in check_title_alignment_consistency(list(prs.slides)).items():
        if sn in slide_reports:
            slide_reports[sn]["issues"].extend(issues)

    for sn, issues in check_title_bold_consistency(list(prs.slides)).items():
        if sn in slide_reports:
            slide_reports[sn]["issues"].extend(issues)

    # Поля по слайду 2
    for sn, issues in check_content_margins(list(prs.slides)).items():
        if sn in slide_reports:
            slide_reports[sn]["issues"].extend(issues)

    if template_path:
        for sn, issues in compare_with_template(template_path, filepath).items():
            if sn in slide_reports:
                slide_reports[sn]["issues"].extend(issues)

    total  = sum(len(v["issues"]) for v in slide_reports.values())
    errors = sum(1 for v in slide_reports.values()
                 for i in v["issues"] if i.get("severity") == "error")
    codes  = Counter(i["code"] for v in slide_reports.values() for i in v["issues"])

    return {
        "total_slides": total_slides,
        "slides": slide_reports,
        "with_template": template_path is not None,
        "summary": {"total": total, "errors": errors, "warnings": total - errors, "codes": codes},
    }


# ─────────────────────────────────────────────
# ФОРМАТИРОВАНИЕ
# ─────────────────────────────────────────────

CODE_LABELS = {
    "HANGING_PREPOSITION":      "Висячий предлог",
    "DOUBLE_SPACE":             "Двойной пробел",
    "HYPHEN_AS_DASH":           "Дефис вместо тире",
    "EN_DASH_AS_EM":            "Короткое тире (–) вместо длинного (—)",
    "RANGE_HYPHEN":             "Дефис в диапазоне",
    "WRONG_QUOTES":             "Прямые кавычки",
    "PERCENT_SPACE":            "Процент с пробелом",
    "NO_SPACE_UNIT":            "Нет пробела перед сокращением",
    "RUB_NO_DOT":               "«руб» без точки",
    "MIN_NO_DOT":               "«мин» без точки",
    "CURRENCY_BEFORE":          "₽ перед числом",
    "MULTIPLY_X":               "Умножение буквой X",
    "TITLE_DOT":                "Точка в заголовке",
    "DASH_NO_SPACE":            "Тире без пробелов",
    "LARGE_NUMBER":             "Большое число без пробелов",
    "MISSING_YO":               "Буква «ё»",
    "TITLE_TOO_LONG":           "Заголовок слишком длинный",
    "FONT_TOO_SMALL":           "Шрифт слишком маленький",
    "TOO_MANY_FONT_SIZES":      "Много разных размеров шрифта",
    "LINE_SPACING_TIGHT":       "Межстрочный интервал плотный",
    "LINE_SPACING_LOOSE":       "Межстрочный интервал широкий",
    "BULLET_INCONSISTENT":      "Разные типы буллитов",
    "BULLET_LEVEL_GAP":         "Пропуск уровня отступа",
    "TITLE_FONT_INCONSISTENT":  "Шрифт заголовка скачет",
    "BODY_FONT_INCONSISTENT":   "Шрифт текста скачет",
    "TITLE_FONT_FAMILY_INCONSISTENT": "Семейство шрифта заголовка отличается",
    "BODY_FONT_FAMILY_INCONSISTENT":  "Семейство шрифта текста отличается",
    "ELEMENT_OVERLAP":          "Элементы перекрываются",
    "NEAR_ALIGNED":             "Почти выровнено",
    "TIGHT_SPACING":            "Слишком маленький отступ",
    "IMAGE_DISTORTED":          "Изображение искажено",
    "IMAGE_SIZE_MISMATCH":      "Размер изображения отличается",
    "IMAGE_POS_MISMATCH":       "Позиция изображения отличается",
    "TOO_MANY_COLORS":          "Много цветов",
    "TABLE_TOO_MANY_ROWS":      "Таблица перегружена строками",
    "TABLE_TOO_MANY_COLS":      "Таблица перегружена столбцами",
    "TABLE_DUPLICATE_HEADERS":  "Повторы в заголовках таблицы",
    "TABLE_NUMBER_ALIGN":       "Числа в таблице по левому краю",
    "SHAPE_TOP_MISMATCH":       "Элемент сдвинут вертикально",
    "SHAPE_LEFT_MISMATCH":      "Элемент сдвинут горизонтально",
    "SHAPE_WIDTH_MISMATCH":     "Ширина изменена",
    "SHAPE_HEIGHT_MISMATCH":    "Высота изменена",
    "SHAPE_FONT_MISMATCH":      "Шрифт не совпадает",
    "FONT_SIZE_NOT_ALLOWED":    "Нестандартный размер шрифта",
    "FONT_MIX":                 "Смешение шрифтовых семейств",
    "TITLE_ALIGN_INCONSISTENT": "Выравнивание заголовка отличается",
    "TITLE_BOLD_INCONSISTENT":  "Начертание заголовка отличается",
    "CONTENT_MARGIN_LEFT":      "Элемент левее поля",
    "LINE_SPACING_GUIDE":       "Межстрочный не по гайду",
    "ROUNDNESS_MISMATCH":       "Скругление фигуры изменено",
    "TMPL_TOP":                 "Смещение от шаблона (вертикаль)",
    "TMPL_LEFT":                "Смещение от шаблона (горизонталь)",
    "TMPL_COUNT":               "Количество элементов ≠ шаблону",
}


def format_report(report: dict) -> str:
    lines = []
    s = report["summary"]
    tmpl = " + шаблон" if report.get("with_template") else ""

    lines.append(f"Отчёт по презентации{tmpl}")
    lines.append(f"{report['total_slides']} слайдов  ·  {s['errors']} ошибки  ·  {s['warnings']} замечания")

    if s["total"] == 0:
        lines.append("\n✅ Ошибок не найдено")
        return "\n".join(lines)

    for slide_num, data in sorted(report["slides"].items()):
        issues = data["issues"]

        # Разделитель между слайдами
        lines.append("")
        lines.append(f"*Слайд {slide_num}* — {data['title']}")
        lines.append("━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━")

        if not issues:
            lines.append("✅ Ошибок нет")
            continue

        errors   = [i for i in issues if i.get("severity") == "error"]
        warnings = [i for i in issues if i.get("severity") == "warning"]

        for i in errors:
            lines.append(f"❌ {i['message']}")
            if i.get("line"):
                snippet = i["line"][:60].replace("*", "\\*").replace("_", "\\_")
                lines.append(f"   ↳ _{snippet}_")
            lines.append("")   # пустая строка после каждого пункта

        for i in warnings:
            lines.append(f"⚠️ {i['message']}")
            if i.get("line"):
                snippet = i["line"][:60].replace("*", "\\*").replace("_", "\\_")
                lines.append(f"   ↳ _{snippet}_")
            lines.append("")   # пустая строка после каждого пункта

    if s["codes"]:
        lines.append("━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━")
        lines.append("*Чаще всего:*")
        for code, count in s["codes"].most_common(5):
            lines.append(f"  {CODE_LABELS.get(code, code)}: {count}×")

    return "\n".join(lines)


def split_message(text: str, limit: int = 4000) -> list:
    parts, current, cur_len = [], [], 0
    for line in text.split("\n"):
        if cur_len + len(line) + 1 > limit and current:
            parts.append("\n".join(current))
            current, cur_len = [], 0
        current.append(line)
        cur_len += len(line) + 1
    if current:
        parts.append("\n".join(current))
    return parts


def extract_all_text(filepath: str) -> str:
    prs = Presentation(filepath)
    result = []
    for idx, slide in enumerate(prs.slides):
        lines = [f"Слайд {idx + 1}:"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = get_shape_text(shape).strip()
                if t:
                    lines.append(t)
        if len(lines) > 1:
            result.append("\n".join(lines))
    return "\n\n".join(result)
